"""
AGRIKONSULT FFS Livestock Value Chain Strategy - PowerPoint Generator
Run this script to generate the PowerPoint presentation.

Requirements:
    pip install python-pptx

Usage:
    python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
import os

# Colors
PRIMARY_GREEN = RgbColor(0x1B, 0x5E, 0x20)
SECONDARY_GREEN = RgbColor(0x2E, 0x7D, 0x32)
ACCENT_GREEN = RgbColor(0x4C, 0xAF, 0x50)
LIGHT_GREEN = RgbColor(0x81, 0xC7, 0x84)
PALE_GREEN = RgbColor(0xE8, 0xF5, 0xE9)
GOLD = RgbColor(0xFF, 0xB3, 0x00)
DARK = RgbColor(0x1A, 0x1A, 0x1A)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
RED = RgbColor(0xE5, 0x39, 0x35)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_shape_gradient(shape, color1, color2):
    """Set gradient fill for a shape"""
    shape.fill.gradient()
    shape.fill.gradient_angle = 135
    shape.fill.gradient_stops[0].color.rgb = color1
    shape.fill.gradient_stops[1].color.rgb = color2

def add_text_box(slide, left, top, width, height, text, font_size=18, font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to a slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return textbox

def add_stat_card(slide, left, top, number, label, width=Inches(2.2), height=Inches(1.5)):
    """Add a statistics card"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GREEN
    card.line.width = Pt(2)

    # Number
    add_text_box(slide, left + Inches(0.1), top + Inches(0.3), width - Inches(0.2), Inches(0.6),
                 number, font_size=32, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # Label
    add_text_box(slide, left + Inches(0.1), top + Inches(0.9), width - Inches(0.2), Inches(0.5),
                 label, font_size=11, font_color=DARK, alignment=PP_ALIGN.CENTER)

def add_pillar_target(slide, left, top, value, label):
    """Add a target metric badge"""
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.9))
    set_shape_fill(badge, PALE_GREEN)
    badge.line.color.rgb = LIGHT_GREEN

    add_text_box(slide, left + Inches(0.1), top + Inches(0.1), Inches(1.6), Inches(0.4),
                 value, font_size=24, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.5), Inches(1.6), Inches(0.35),
                 label, font_size=9, font_color=DARK, alignment=PP_ALIGN.CENTER)

def create_title_slide(prs):
    """Slide 1: Title"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Green gradient background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(background, PRIMARY_GREEN, ACCENT_GREEN)
    background.line.fill.background()

    # Logo placeholder
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(0.4), Inches(1.6), Inches(0.7))
    set_shape_fill(logo_bg, WHITE)
    logo_bg.line.fill.background()
    add_text_box(slide, Inches(4.3), Inches(0.45), Inches(1.4), Inches(0.3),
                 "AGRIKONSULT", font_size=12, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.3), Inches(0.7), Inches(1.4), Inches(0.25),
                 "Farm Solutions Nigeria Ltd", font_size=7, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(1.5),
                 "Livestock Value Chain Strategy\nfor the Farmers' Field School Programme",
                 font_size=40, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")

    # Subtitle
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.6),
                 "Transforming Nigeria's Livestock Sector through Knowledge-Driven, Market-Oriented Solutions",
                 font_size=18, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # Meta info
    add_text_box(slide, Inches(0.5), Inches(5), Inches(9), Inches(0.8),
                 "Policy Workshop & Stakeholder Briefing\nJanuary 2026",
                 font_size=14, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # Slide number
    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "1 / 20", font_size=10, font_color=LIGHT_GREEN, alignment=PP_ALIGN.RIGHT)

def create_stats_slide(prs):
    """Slide 2: Nigeria's Livestock Sector Statistics"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    # Title
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(6), Inches(0.7))
    title_shape.fill.background()
    title_shape.line.fill.background()
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(8), Inches(0.7),
                 "Nigeria's Livestock Sector: The Opportunity",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    # Stats cards - Row 1
    add_stat_card(slide, Inches(0.5), Inches(1.6), "563M", "Chickens")
    add_stat_card(slide, Inches(2.9), Inches(1.6), "124M", "Goats")
    add_stat_card(slide, Inches(5.3), Inches(1.6), "60M", "Sheep")

    # Stats cards - Row 2
    add_stat_card(slide, Inches(0.5), Inches(3.3), "58M", "Cattle")
    add_stat_card(slide, Inches(2.9), Inches(3.3), "17%", "of Agricultural GDP")
    add_stat_card(slide, Inches(5.3), Inches(3.3), "25%", "of West Africa's Livestock")

    # Highlight box
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(7), Inches(0.7))
    set_shape_fill(highlight, PALE_GREEN)
    highlight.line.color.rgb = ACCENT_GREEN
    highlight.line.width = Pt(3)

    add_text_box(slide, Inches(0.7), Inches(5.2), Inches(6.6), Inches(0.5),
                 "Nigeria is the leading livestock producer in West Africa, with untapped potential of ₦30 trillion",
                 font_size=14, font_color=PRIMARY_GREEN, bold=False, alignment=PP_ALIGN.CENTER)

    # Slide number
    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "2 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_vision_slide(prs):
    """Slide 3: The Vision $74 Billion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "The National Vision: $74 Billion by 2035",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    # Left column text
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4),
                 "Current State → Future Target", font_size=18, font_color=PRIMARY_GREEN, bold=True)

    add_text_box(slide, Inches(0.5), Inches(2), Inches(4.5), Inches(1.2),
                 "The Federal Ministry of Livestock Development has set an ambitious goal to grow the sector from $32 billion to $74 billion contribution to GDP by 2035.\n\nThis represents a 131% increase in just one decade.",
                 font_size=12, font_color=DARK)

    # Highlight box
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(4.5), Inches(0.7))
    set_shape_fill(highlight, PALE_GREEN)
    highlight.line.color.rgb = ACCENT_GREEN
    add_text_box(slide, Inches(0.6), Inches(3.9), Inches(4.3), Inches(0.5),
                 "$2.5 Billion JBS Investment signed in November 2024 for 6 processing plants",
                 font_size=11, font_color=PRIMARY_GREEN, bold=True)

    # Right column - stat cards
    add_stat_card(slide, Inches(5.5), Inches(1.5), "$32B", "Current Contribution (2024)", width=Inches(2.3))
    add_stat_card(slide, Inches(5.5), Inches(3.2), "$74B", "Target Contribution (2035)", width=Inches(2.3))
    add_stat_card(slide, Inches(5.5), Inches(4.9), "2M+", "Jobs by 2028 (NLTP Target)", width=Inches(2.3))

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "3 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_challenges_slide(prs):
    """Slide 4: Current Challenges"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Current Challenges Facing the Sector",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    challenges = [
        ("📉", "Low productivity and high mortality rates"),
        ("💊", "Weak input and service delivery systems"),
        ("🏠", "Poor husbandry and biosecurity practices"),
        ("💰", "Limited access to finance and markets"),
        ("🥩", "High post-production losses and food safety risks"),
        ("🌡️", "Climate and disease shocks (ND, ASF, PPR, CBPP)")
    ]

    for i, (icon, text) in enumerate(challenges):
        row = i // 2
        col = i % 2
        left = Inches(0.5 + col * 4.3)
        top = Inches(1.5 + row * 1.1)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(0.9))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = RED
        card.line.width = Pt(2)

        add_text_box(slide, left + Inches(0.15), top + Inches(0.25), Inches(3.7), Inches(0.5),
                     f"{icon}  {text}", font_size=12, font_color=DARK)

    # Quote box
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(7.5), Inches(1))
    set_shape_fill(quote_box, WHITE)
    quote_box.line.color.rgb = LIGHT_GREEN

    add_text_box(slide, Inches(0.7), Inches(5.15), Inches(7.1), Inches(0.7),
                 '"Per capita beef consumption in Nigeria is just 1.7 kg/year compared to the global average of 9.1 kg — a clear indicator of unmet demand and opportunity."',
                 font_size=12, font_color=DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "4 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_ffs_solution_slide(prs):
    """Slide 5: The FFS Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "The Farmers' Field School Solution",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    # Left content
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4),
                 "What is FFS?", font_size=18, font_color=PRIMARY_GREEN, bold=True)

    add_text_box(slide, Inches(0.5), Inches(1.9), Inches(4.5), Inches(1.2),
                 "A Farmers' Field School brings together groups of farmers, livestock herders, or fisherfolk to learn how to shift towards more sustainable, productive practices through experiential learning and collective action.",
                 font_size=11, font_color=DARK)

    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(4.5), Inches(0.4),
                 "Why FFS for Livestock?", font_size=18, font_color=PRIMARY_GREEN, bold=True)

    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(4.5), Inches(1),
                 "The FFS model provides a powerful platform to address constraints across the entire value chain — not just on-farm production — transforming subsistence farmers into market-oriented entrepreneurs.",
                 font_size=11, font_color=DARK)

    # Right stats
    add_stat_card(slide, Inches(5.5), Inches(1.5), "30+", "Years of Proven FFS Success", width=Inches(2.3))
    add_stat_card(slide, Inches(5.5), Inches(3.2), "2,000+", "Active FFS in Eastern Africa", width=Inches(2.3))
    add_stat_card(slide, Inches(5.5), Inches(4.9), "50,000+", "Direct Beneficiaries", width=Inches(2.3))

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "5 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_objectives_slide(prs):
    """Slide 6: Strategic Objectives"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Strategic Objectives",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(3.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    objectives = [
        "ICT-Innovative Approach — Knowledge-driven livestock extension delivery",
        "Improve Productivity — Through Good Animal Husbandry Practices (GAHP)",
        "Strengthen Capacity — Meet market, quality, and food safety standards",
        "Enhance Value Addition — Aggregation and collective marketing",
        "Improve Access — Inputs, services, finance, and insurance",
        "Climate-Smart Systems — Gender-responsive, youth-inclusive",
        "Strengthen Linkages — Between farmers, processors, traders, regulators"
    ]

    for i, obj in enumerate(objectives):
        row = i // 2
        col = i % 2
        if i == 6:  # Last item centered
            left = Inches(2.4)
        else:
            left = Inches(0.5 + col * 4.3)
        top = Inches(1.4 + row * 0.95)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(0.8))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = LIGHT_GREEN

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1), top + Inches(0.2), Inches(0.4), Inches(0.4))
        set_shape_gradient(circle, PRIMARY_GREEN, ACCENT_GREEN)
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.1), top + Inches(0.25), Inches(0.4), Inches(0.35),
                     str(i + 1), font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.6), top + Inches(0.2), Inches(3.2), Inches(0.5),
                     obj, font_size=10, font_color=DARK)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "6 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_framework_slide(prs):
    """Slide 7: Value Chain Framework"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Livestock Value Chain Framework within FFS",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(0.4),
                 "Seven interconnected nodes translated into FFS learning modules, field experiments, and collective business actions",
                 font_size=12, font_color=DARK, alignment=PP_ALIGN.CENTER)

    nodes = [
        ("1", "Research-Extension\nLinkage"),
        ("2", "Input Supply\nSystems"),
        ("3", "Production &\nHusbandry"),
        ("4", "Animal Health\n& Biosecurity"),
        ("5", "Aggregation &\nProcessing"),
        ("6", "Marketing &\nMarket Systems"),
        ("7", "Support Services")
    ]

    # Row 1 - nodes 1-4
    for i in range(4):
        left = Inches(0.4 + i * 2.2)
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(1.8), Inches(1.1))
        set_shape_gradient(node, PRIMARY_GREEN, SECONDARY_GREEN)
        node.line.fill.background()

        add_text_box(slide, left, Inches(2.25), Inches(1.8), Inches(0.4),
                     nodes[i][0], font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(2.65), Inches(1.8), Inches(0.6),
                     nodes[i][1], font_size=9, font_color=WHITE, alignment=PP_ALIGN.CENTER)

        if i < 3:
            add_text_box(slide, left + Inches(1.85), Inches(2.55), Inches(0.3), Inches(0.3),
                         "→", font_size=18, font_color=ACCENT_GREEN, bold=True)

    # Row 2 - nodes 5-7
    for i in range(3):
        left = Inches(1.5 + i * 2.4)
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4), Inches(1.8), Inches(1.1))
        set_shape_gradient(node, PRIMARY_GREEN, SECONDARY_GREEN)
        node.line.fill.background()

        add_text_box(slide, left, Inches(4.05), Inches(1.8), Inches(0.4),
                     nodes[i + 4][0], font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.45), Inches(1.8), Inches(0.6),
                     nodes[i + 4][1], font_size=9, font_color=WHITE, alignment=PP_ALIGN.CENTER)

        if i < 2:
            add_text_box(slide, left + Inches(1.85), Inches(4.35), Inches(0.3), Inches(0.3),
                         "→", font_size=18, font_color=ACCENT_GREEN, bold=True)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "7 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_pillar_slide(prs, slide_num, pillar_num, title, subtitle, issues, strategies, targets):
    """Create a strategic pillar slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
                 f"Strategic Pillar {pillar_num}: {title}",
                 font_size=24, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(5), Inches(0.04))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    # Pillar card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(7.5), Inches(4.2))
    set_shape_fill(card, WHITE)
    card.line.color.rgb = LIGHT_GREEN

    # Pillar number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(0.5), Inches(0.5))
    set_shape_gradient(badge, PRIMARY_GREEN, ACCENT_GREEN)
    badge.line.fill.background()
    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(0.5), Inches(0.45),
                 str(pillar_num), font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.3), Inches(1.4), Inches(6), Inches(0.4),
                 subtitle, font_size=14, font_color=PRIMARY_GREEN, bold=True)

    # Issues column
    add_text_box(slide, Inches(0.7), Inches(2), Inches(3.3), Inches(0.3),
                 "🔍 Key Issues", font_size=11, font_color=SECONDARY_GREEN, bold=True)

    issues_text = "\n".join([f"• {issue}" for issue in issues])
    add_text_box(slide, Inches(0.7), Inches(2.3), Inches(3.3), Inches(1.5),
                 issues_text, font_size=10, font_color=DARK)

    # Strategies column
    add_text_box(slide, Inches(4.2), Inches(2), Inches(3.5), Inches(0.3),
                 "📋 FFS Strategies", font_size=11, font_color=SECONDARY_GREEN, bold=True)

    strategies_text = "\n".join([f"• {strat}" for strat in strategies])
    add_text_box(slide, Inches(4.2), Inches(2.3), Inches(3.5), Inches(1.8),
                 strategies_text, font_size=10, font_color=DARK)

    # Target metrics
    for i, (value, label) in enumerate(targets):
        add_pillar_target(slide, Inches(0.7 + i * 2.2), Inches(4.4), value, label)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 f"{slide_num} / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_cross_cutting_slide(prs):
    """Slide 14: Cross-Cutting Themes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Cross-Cutting Themes",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(3.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    themes = [
        ("👥", "Gender & Social Inclusion", ["Targeted women-only and youth FFS groups", "Enterprises with low entry barriers", "Leadership and entrepreneurship training"]),
        ("🌱", "Climate-Smart Practices", ["Heat stress management and water efficiency", "Improved fodder, hay, silage systems", "Manure management and biogas"]),
        ("🥗", "Nutrition & Food Safety", ["Animal-source foods for household nutrition", "Hygiene and slaughter safety practices", "Milk handling and safety protocols"])
    ]

    for i, (icon, title, items) in enumerate(themes):
        left = Inches(0.4 + i * 2.9)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(2.7), Inches(3.5))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = LIGHT_GREEN

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), Inches(1.7), Inches(0.8), Inches(0.8))
        set_shape_gradient(circle, PRIMARY_GREEN, ACCENT_GREEN)
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.95), Inches(1.85), Inches(0.8), Inches(0.5),
                     icon, font_size=24, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.1), Inches(2.6), Inches(2.5), Inches(0.4),
                     title, font_size=12, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

        items_text = "\n".join([f"• {item}" for item in items])
        add_text_box(slide, left + Inches(0.15), Inches(3.1), Inches(2.4), Inches(1.7),
                     items_text, font_size=9, font_color=DARK)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "14 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_alignment_slide(prs):
    """Slide 15: Alignment with National Policies"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                 "Alignment with National Policies",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    # L-PRES Card
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.4), Inches(4), Inches(3.3))
    set_shape_fill(card1, WHITE)
    card1.line.color.rgb = LIGHT_GREEN

    badge1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.55), Inches(0.9), Inches(0.35))
    set_shape_gradient(badge1, PRIMARY_GREEN, ACCENT_GREEN)
    badge1.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(1.57), Inches(0.9), Inches(0.3),
                 "L-PRES", font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.6), Inches(1.55), Inches(2.5), Inches(0.35),
                 "$500M World Bank Project", font_size=11, font_color=DARK)

    lpres_items = "✓ 1.4 million beneficiaries targeted\n✓ Focus on beef cattle, sheep & goat\n✓ Crisis prevention and conflict mitigation\n✓ Institutional strengthening\n✓ Climate adaptation with NiMet"
    add_text_box(slide, Inches(0.6), Inches(2.1), Inches(3.6), Inches(2.4),
                 lpres_items, font_size=10, font_color=DARK)

    # NLTP Card
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.4), Inches(4), Inches(3.3))
    set_shape_fill(card2, WHITE)
    card2.line.color.rgb = LIGHT_GREEN

    badge2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.55), Inches(0.8), Inches(0.35))
    set_shape_gradient(badge2, PRIMARY_GREEN, ACCENT_GREEN)
    badge2.line.fill.background()
    add_text_box(slide, Inches(4.8), Inches(1.57), Inches(0.8), Inches(0.3),
                 "NLTP", font_size=10, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(5.7), Inches(1.55), Inches(2.7), Inches(0.35),
                 "National Transformation Plan", font_size=11, font_color=DARK)

    nltp_items = "✓ ₦100 billion budget (80% Fed, 20% States)\n✓ 2 million+ jobs target by 2028\n✓ Six pillars for transformation\n✓ Modernizing through ranching\n✓ Farmer-herder conflict resolution"
    add_text_box(slide, Inches(4.8), Inches(2.1), Inches(3.6), Inches(2.4),
                 nltp_items, font_size=10, font_color=DARK)

    # Bottom highlight
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.9), Inches(8.2), Inches(0.65))
    set_shape_fill(highlight, PALE_GREEN)
    highlight.line.color.rgb = ACCENT_GREEN

    add_text_box(slide, Inches(0.5), Inches(5), Inches(8), Inches(0.5),
                 "This FFS Livestock Strategy directly supports both L-PRES and NLTP objectives through grassroots capacity building",
                 font_size=11, font_color=PRIMARY_GREEN, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "15 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_success_stories_slide(prs):
    """Slide 16: FFS Success Stories"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 "FFS Success Stories from Africa",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    stories = [
        ("🇹🇿 Tanzania", "FAO integrated behavioural science into FFS to address antimicrobial resistance in broiler production.", "Scaling up across the region"),
        ("🇿🇼 Zimbabwe", "FAO piloted FFS in 8 districts, graduating 106 FFS farmers and facilitators.", "106 graduates as multipliers"),
        ("🇰🇪 Kenya", "Agropastoralist field schools built drought resilience through fodder production.", "FFS in national extension policy"),
        ("🇺🇬 Uganda", "Studies demonstrated clear benefits in milk production and dairy management.", "Proven dairy improvements")
    ]

    for i, (country, desc, stat) in enumerate(stories):
        row = i // 2
        col = i % 2
        left = Inches(0.4 + col * 4.3)
        top = Inches(1.3 + row * 1.7)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.1), Inches(1.5))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = ACCENT_GREEN
        card.line.width = Pt(2)

        add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.8), Inches(0.35),
                     country, font_size=12, font_color=PRIMARY_GREEN, bold=True)

        add_text_box(slide, left + Inches(0.15), top + Inches(0.45), Inches(3.8), Inches(0.65),
                     desc, font_size=9, font_color=DARK)

        stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(1.1), Inches(3.8), Inches(0.3))
        set_shape_fill(stat_box, PALE_GREEN)
        stat_box.line.fill.background()
        add_text_box(slide, left + Inches(0.2), top + Inches(1.12), Inches(3.7), Inches(0.25),
                     stat, font_size=9, font_color=PRIMARY_GREEN, bold=True)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "16 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_implementation_slide(prs):
    """Slide 17: Implementation Arrangements"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 "Implementation Arrangements",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(4), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    partners = [
        ("🏛️", "Lead Institutions", "FMLD & State Ministries"),
        ("🔬", "Technical Partners", "Research, Universities, NGOs"),
        ("👨‍🏫", "Delivery Agents", "Extension Officers as Facilitators"),
        ("🏢", "Private Sector", "Input Suppliers, Processors")
    ]

    for i, (icon, ptype, pname) in enumerate(partners):
        left = Inches(0.4 + i * 2.15)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.3), Inches(2), Inches(1.4))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = LIGHT_GREEN

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), Inches(1.4), Inches(0.7), Inches(0.7))
        set_shape_fill(circle, PALE_GREEN)
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.6), Inches(1.5), Inches(0.7), Inches(0.5),
                     icon, font_size=20, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.1), Inches(2.15), Inches(1.8), Inches(0.3),
                     ptype, font_size=9, font_color=SECONDARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(2.4), Inches(1.8), Inches(0.3),
                     pname, font_size=8, font_color=DARK, alignment=PP_ALIGN.CENTER)

    # Timeline
    add_text_box(slide, Inches(0.5), Inches(3), Inches(8), Inches(0.3),
                 "Implementation Timeline", font_size=14, font_color=PRIMARY_GREEN, bold=True)

    # Timeline line
    timeline_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.7), Inches(7), Inches(0.05))
    set_shape_fill(timeline_line, LIGHT_GREEN)
    timeline_line.line.fill.background()

    phases = [
        ("1", "Months 1-3", "Facilitator Training"),
        ("2", "Months 3-6", "FFS Establishment"),
        ("3", "Months 6-12", "Learning Cycles"),
        ("4", "Month 12+", "Graduation & Scaling")
    ]

    for i, (num, period, desc) in enumerate(phases):
        left = Inches(0.8 + i * 2.1)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.35), Inches(3.45), Inches(0.5), Inches(0.5))
        set_shape_gradient(circle, PRIMARY_GREEN, ACCENT_GREEN)
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.35), Inches(3.5), Inches(0.5), Inches(0.4),
                     num, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left, Inches(4.05), Inches(1.2), Inches(0.25),
                     period, font_size=9, font_color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.3), Inches(1.2), Inches(0.25),
                     desc, font_size=8, font_color=DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(8), Inches(0.3),
                 "FFS cycles run for 9-12 months, aligned with livestock production cycles",
                 font_size=11, font_color=DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "17 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_mel_slide(prs):
    """Slide 18: Monitoring, Evaluation & Learning"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 "Monitoring, Evaluation & Learning",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(8), Inches(0.4),
                 "Participatory M&E tools embedded within FFS sessions for continuous improvement",
                 font_size=12, font_color=DARK)

    indicators = [
        ("📊", "Productivity Metrics", "Mortality rates, growth rates, output per animal"),
        ("💵", "Economic Indicators", "Income levels, profitability, cost reduction"),
        ("✅", "Practice Adoption", "GAHP and biosecurity practice uptake rates"),
        ("🛒", "Market Participation", "Collective sales, market linkages, price realization")
    ]

    for i, (icon, title, desc) in enumerate(indicators):
        row = i // 2
        col = i % 2
        left = Inches(0.4 + col * 4.3)
        top = Inches(1.9 + row * 1.3)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(1.1))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = LIGHT_GREEN

        icon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.2), Inches(0.6), Inches(0.6))
        set_shape_gradient(icon_box, PRIMARY_GREEN, ACCENT_GREEN)
        icon_box.line.fill.background()
        add_text_box(slide, left + Inches(0.15), top + Inches(0.3), Inches(0.6), Inches(0.4),
                     icon, font_size=18, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.9), top + Inches(0.2), Inches(2.9), Inches(0.35),
                     title, font_size=11, font_color=DARK, bold=True)
        add_text_box(slide, left + Inches(0.9), top + Inches(0.55), Inches(2.9), Inches(0.45),
                     desc, font_size=9, font_color=DARK)

    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.7), Inches(8.2), Inches(0.6))
    set_shape_fill(highlight, PALE_GREEN)
    highlight.line.color.rgb = ACCENT_GREEN

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(8), Inches(0.4),
                 "Data-driven decision making ensures continuous improvement and evidence-based scaling",
                 font_size=11, font_color=PRIMARY_GREEN, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "18 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_scaling_slide(prs):
    """Slide 19: Scaling & Sustainability"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, WHITE, PALE_GREEN)
    bg.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 "Scaling & Sustainability Strategy",
                 font_size=28, font_color=PRIMARY_GREEN, bold=True, font_name="Georgia")

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(4.5), Inches(0.05))
    set_shape_fill(line, ACCENT_GREEN)
    line.line.fill.background()

    strategies = [
        ("🎓", "FFS Graduation", "Groups transition into cooperatives or agribusiness clusters"),
        ("👥", "Community Facilitators", "Training of CBFs for local ownership"),
        ("📋", "Policy Integration", "Integration into state livestock development plans"),
        ("🤝", "Blended Finance", "Public-Private Partnership models for sustainable funding")
    ]

    for i, (icon, title, desc) in enumerate(strategies):
        left = Inches(0.4 + i * 2.15)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.4), Inches(2), Inches(1.8))
        set_shape_fill(card, WHITE)
        card.line.color.rgb = LIGHT_GREEN

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.6), Inches(1.55), Inches(0.7), Inches(0.7))
        set_shape_gradient(circle, PRIMARY_GREEN, ACCENT_GREEN)
        circle.line.fill.background()
        add_text_box(slide, left + Inches(0.6), Inches(1.65), Inches(0.7), Inches(0.5),
                     icon, font_size=20, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.1), Inches(2.35), Inches(1.8), Inches(0.35),
                     title, font_size=10, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), Inches(2.7), Inches(1.8), Inches(0.45),
                     desc, font_size=8, font_color=DARK, alignment=PP_ALIGN.CENTER)

    # Quote box
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.5), Inches(8.2), Inches(1.2))
    set_shape_fill(quote, WHITE)
    quote.line.color.rgb = LIGHT_GREEN

    add_text_box(slide, Inches(0.6), Inches(3.6), Inches(7.8), Inches(1),
                 '"FFS networks in Eastern Africa have clearly shown how farmers themselves have been able to build bottom-up producer organizations during and after projects ended — a model Nigeria can replicate."',
                 font_size=11, font_color=DARK, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "19 / 20", font_size=10, font_color=SECONDARY_GREEN, alignment=PP_ALIGN.RIGHT)

def create_contact_slide(prs):
    """Slide 20: Conclusion & Contact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_gradient(bg, PRIMARY_GREEN, SECONDARY_GREEN)
    bg.line.fill.background()

    # Contact card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RgbColor(0xFF, 0xFF, 0xFF)
    card.fill.fore_color.brightness = 0.85
    card.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(0.7),
                 "Transforming Livestock,\nTransforming Lives",
                 font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Georgia")

    add_text_box(slide, Inches(1.8), Inches(2.8), Inches(6.4), Inches(1),
                 "Integrating a livestock value chain strategy into Nigeria's FFS program offers a powerful pathway to transform smallholder livestock systems into resilient, market-linked, and inclusive value chains.",
                 font_size=12, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4), Inches(4), Inches(0.02))
    set_shape_fill(divider, LIGHT_GREEN)
    divider.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(7), Inches(0.4),
                 "Dr. SOSINA, A.O. (Ph.D., RAS)",
                 font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.7), Inches(7), Inches(0.35),
                 "Lead Strategist",
                 font_size=12, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # Email button
    email_btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(5.2), Inches(3.6), Inches(0.5))
    set_shape_fill(email_btn, WHITE)
    email_btn.line.fill.background()
    add_text_box(slide, Inches(3.2), Inches(5.25), Inches(3.6), Inches(0.4),
                 "dayososina@gmail.com",
                 font_size=12, font_color=PRIMARY_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(9), Inches(7), Inches(0.8), Inches(0.3),
                 "20 / 20", font_size=10, font_color=LIGHT_GREEN, alignment=PP_ALIGN.RIGHT)

def main():
    """Generate the complete presentation"""
    print("Creating AGRIKONSULT FFS Presentation...")

    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(prs)
    print("  ✓ Slide 1: Title")

    # Slide 2: Statistics
    create_stats_slide(prs)
    print("  ✓ Slide 2: Statistics")

    # Slide 3: Vision
    create_vision_slide(prs)
    print("  ✓ Slide 3: Vision")

    # Slide 4: Challenges
    create_challenges_slide(prs)
    print("  ✓ Slide 4: Challenges")

    # Slide 5: FFS Solution
    create_ffs_solution_slide(prs)
    print("  ✓ Slide 5: FFS Solution")

    # Slide 6: Objectives
    create_objectives_slide(prs)
    print("  ✓ Slide 6: Objectives")

    # Slide 7: Framework
    create_framework_slide(prs)
    print("  ✓ Slide 7: Framework")

    # Slides 8-13: Strategic Pillars
    pillars_data = [
        (8, 1, "Input Supply & Service Systems", "Addressing Poor Quality Feed, Drugs, and Breeding Stock",
         ["Poor quality feed, drugs, breeding stock", "Limited veterinary/extension services", "Lack of standards and quality control"],
         ["Input quality testing sessions", "Local feed formulation", "Community-based bulk purchasing", "Vet and agro-dealer engagement"],
         [("20%", "Cost Reduction"), ("10%", "Quality Improvement")]),

        (9, 2, "Production & GAHP", "Species-Specific Husbandry Excellence",
         ["Poultry, goats, sheep, cattle, pigs", "Housing, feeding, watering standards", "Animal welfare practices"],
         ["Season-long learning cycles", "Comparative trials", "Farmer-led data collection"],
         [("20%", "Productivity Increase"), ("10%", "Mortality Reduction"), ("30%", "Standardized Practices")]),

        (10, 3, "Animal Health & Biosecurity", "Combating Endemic Diseases & Zoonotic Risks",
         ["Endemic diseases: ND, ASF, PPR, CBPP", "Poor biosecurity practices", "Antimicrobial resistance risks"],
         ["Participatory disease surveillance", "Biosecurity mapping", "Vaccination calendars", "AMR stewardship training", "One Health integration"],
         [("20%", "Disease Reduction"), ("20%", "Food Safety Improvement")]),

        (11, 4, "Aggregation & Value Addition", "From Distress Sales to Premium Markets",
         ["Distress sales, low farm-gate prices", "Minimal value addition", "High post-harvest losses"],
         ["Collective aggregation models", "Basic processing training", "Women/youth micro-processing", "Cost-benefit analysis"],
         [("30%", "Margin Increase"), ("20%", "Loss Reduction")]),

        (12, 5, "Marketing & Business Skills", "Market Access & Commercial Orientation",
         ["Weak market information", "Poor bargaining power", "Limited digital adoption"],
         ["Market mapping and price discovery", "Standards and grading training", "Collective marketing simulations", "Digital tools adoption"],
         [("20%", "Market Access"), ("10%", "Commercial Transition")]),

        (13, 6, "Finance & Institutions", "Access to Finance, Insurance & Policy Support",
         ["Limited finance access", "Weak policy implementation", "Lack of institutional support"],
         ["Financial literacy modules", "Microfinance linkages", "Livestock insurance intro", "Government engagement"],
         [("20%", "Investment Increase"), ("10%", "Institutional Support")])
    ]

    for slide_num, pillar_num, title, subtitle, issues, strategies, targets in pillars_data:
        create_pillar_slide(prs, slide_num, pillar_num, title, subtitle, issues, strategies, targets)
        print(f"  ✓ Slide {slide_num}: Pillar {pillar_num}")

    # Slide 14: Cross-Cutting Themes
    create_cross_cutting_slide(prs)
    print("  ✓ Slide 14: Cross-Cutting Themes")

    # Slide 15: Alignment
    create_alignment_slide(prs)
    print("  ✓ Slide 15: Policy Alignment")

    # Slide 16: Success Stories
    create_success_stories_slide(prs)
    print("  ✓ Slide 16: Success Stories")

    # Slide 17: Implementation
    create_implementation_slide(prs)
    print("  ✓ Slide 17: Implementation")

    # Slide 18: MEL
    create_mel_slide(prs)
    print("  ✓ Slide 18: MEL")

    # Slide 19: Scaling
    create_scaling_slide(prs)
    print("  ✓ Slide 19: Scaling")

    # Slide 20: Contact
    create_contact_slide(prs)
    print("  ✓ Slide 20: Contact")

    # Save
    output_path = os.path.join(os.path.dirname(__file__), 'AGRIKONSULT_FFS_Presentation.pptx')
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")
    print("\nNext steps:")
    print("  1. Open the presentation in PowerPoint")
    print("  2. Add animations using: Animations tab → Add Animation")
    print("  3. Recommended: Use 'Fade' or 'Fly In' for content elements")
    print("  4. Set timing: Start 'On Click' or 'After Previous'")

if __name__ == "__main__":
    main()
